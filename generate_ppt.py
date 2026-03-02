"""
SOC Copilot - Project Presentation Generator
Generates an 18-slide professional PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT_BLUE   = RGBColor(0x00, 0x7B, 0xFF)   # Electric blue
ACCENT_CYAN   = RGBColor(0x00, 0xD4, 0xFF)   # Cyan accent
ACCENT_GREEN  = RGBColor(0x00, 0xE6, 0x76)   # Green accent
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)   # Orange accent
ACCENT_RED    = RGBColor(0xFF, 0x3B, 0x30)   # Red accent
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xB0, 0xB8, 0xC8)
MID_GRAY      = RGBColor(0x6C, 0x75, 0x89)
CARD_BG       = RGBColor(0x1A, 0x25, 0x3C)   # Slightly lighter navy


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left, top, width, height, color):
    """Add a colored accent bar/rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_CYAN,
                    font_name="Calibri", spacing=Pt(6)):
    """Add a bulleted list with custom formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.findall(qn('a:buNone'))
        for bn in buNone:
            pPr.remove(bn)
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '▸')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', f'{bullet_color.red:02X}{bullet_color.green:02X}{bullet_color.blue:02X}' if hasattr(bullet_color, 'red') else '00D4FF')
    return txBox


def add_sub_bullets(slide, left, top, width, height, sections,
                    font_size=14, header_color=ACCENT_CYAN, item_color=LIGHT_GRAY):
    """Add sections with sub-bullets. sections = [(header, [items]), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    from lxml import etree

    first = True
    for header, items in sections:
        # Header
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = header
        p.font.size = Pt(font_size + 2)
        p.font.color.rgb = header_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.space_before = Pt(8)

        # Items
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(font_size)
            p.font.color.rgb = item_color
            p.font.name = "Calibri"
            p.space_after = Pt(3)
            p.level = 1
            pPr = p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '•')
            buClr = etree.SubElement(pPr, qn('a:buClr'))
            srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
            srgbClr.set('val', '6C7589')
    return txBox


def add_slide_number(slide, num, total=18):
    """Add slide number at bottom-right."""
    add_text_box(slide, Inches(8.5), Inches(7.1), Inches(1.2), Inches(0.3),
                 f"{num}/{total}", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_header_with_bar(slide, title, slide_num, subtitle=None):
    """Standard slide header with accent bar and title."""
    # Top accent line
    add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), ACCENT_BLUE)
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 title, font_size=32, color=WHITE, bold=True)
    # Underline
    add_accent_bar(slide, Inches(0.8), Inches(1.05), Inches(2), Inches(0.04), ACCENT_CYAN)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.15), Inches(8.4), Inches(0.4),
                     subtitle, font_size=14, color=MID_GRAY)
    add_slide_number(slide, slide_num)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Decorative elements
    add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.08), ACCENT_BLUE)
    add_accent_bar(slide, Inches(0), Inches(7.42), Inches(10), Inches(0.08), ACCENT_CYAN)

    # Left accent stripe
    add_accent_bar(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(1.8), ACCENT_CYAN)

    # Shield icon placeholder
    add_text_box(slide, Inches(0.9), Inches(1.5), Inches(1), Inches(0.8),
                 "🛡️", font_size=40, alignment=PP_ALIGN.LEFT)

    # Main title
    add_text_box(slide, Inches(0.9), Inches(2.0), Inches(8.2), Inches(1.0),
                 "SOC Copilot", font_size=48, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.9), Inches(2.8), Inches(8.2), Inches(0.8),
                 "An AI-Powered Security Operations Assistant for\nAutomated Threat Detection & Incident Response",
                 font_size=20, color=ACCENT_CYAN)

    # Divider
    add_accent_bar(slide, Inches(0.9), Inches(3.8), Inches(3), Inches(0.03), MID_GRAY)

    # Info section
    info_items = [
        ("Institution", "Your Institution Name"),
        ("Department", "Department of Computer Science & Engineering"),
        ("Guide", "Prof. [Guide Name]"),
        ("Team", "[Member 1]  •  [Member 2]  •  [Member 3]  •  [Member 4]"),
        ("Academic Year", "2025 – 2026"),
    ]
    y = 4.1
    for label, value in info_items:
        add_text_box(slide, Inches(0.9), Inches(y), Inches(1.8), Inches(0.3),
                     label, font_size=12, color=MID_GRAY, bold=True)
        add_text_box(slide, Inches(2.7), Inches(y), Inches(6.0), Inches(0.3),
                     value, font_size=12, color=LIGHT_GRAY)
        y += 0.35
    add_slide_number(slide, 1)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 2: Agenda
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Project Overview", 2, "Agenda & Content Flow")

    agenda_left = [
        "Abstract",
        "Problem Statement",
        "Existing System",
        "Limitations",
        "Proposed System",
    ]
    agenda_right = [
        "System Architecture",
        "Data Flow & UML Diagrams",
        "Core & Support Modules",
        "Output Screens",
        "Conclusion & Future Scope",
    ]

    # Left column card
    card = add_accent_bar(slide, Inches(0.8), Inches(1.6), Inches(4.0), Inches(4.8), CARD_BG)
    add_text_box(slide, Inches(1.0), Inches(1.7), Inches(3.6), Inches(0.4),
                 "PART I — Analysis", font_size=14, color=ACCENT_CYAN, bold=True)
    for i, item in enumerate(agenda_left):
        num = f"0{i+1}"
        add_text_box(slide, Inches(1.0), Inches(2.2 + i * 0.65), Inches(0.5), Inches(0.4),
                     num, font_size=24, color=ACCENT_BLUE, bold=True)
        add_text_box(slide, Inches(1.6), Inches(2.25 + i * 0.65), Inches(3.0), Inches(0.4),
                     item, font_size=16, color=LIGHT_GRAY)

    # Right column card
    card = add_accent_bar(slide, Inches(5.2), Inches(1.6), Inches(4.0), Inches(4.8), CARD_BG)
    add_text_box(slide, Inches(5.4), Inches(1.7), Inches(3.6), Inches(0.4),
                 "PART II — Design & Implementation", font_size=14, color=ACCENT_GREEN, bold=True)
    for i, item in enumerate(agenda_right):
        num = f"0{i+6}"
        add_text_box(slide, Inches(5.4), Inches(2.2 + i * 0.65), Inches(0.5), Inches(0.4),
                     num, font_size=24, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, Inches(6.0), Inches(2.25 + i * 0.65), Inches(3.0), Inches(0.4),
                     item, font_size=16, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 3: Abstract
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Abstract", 3)

    abstract_points = [
        "Modern organizations face an ever-increasing volume of cyber threats, making manual security monitoring unsustainable",
        "SOC (Security Operations Center) teams are overwhelmed with thousands of alerts daily, leading to alert fatigue and missed threats",
        "Traditional log analysis is time-consuming and relies heavily on manual investigation by skilled analysts",
        "SOC Copilot is an AI-powered desktop assistant that automates log analysis using hybrid ML — Isolation Forest for anomaly detection and Random Forest for attack classification",
        "Provides intelligent threat detection, natural-language explanations, and automated incident response recommendations — all running fully offline",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.0),
                    abstract_points, font_size=17, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 4: Problem Statement
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Problem Statement", 4)

    problems = [
        "Organizations generate massive volumes of security logs (millions of events/day) that exceed human processing capacity",
        "Alert fatigue — analysts are desensitized to constant warnings, causing critical threats to be overlooked",
        "Slow incident response due to manual log correlation and investigation workflows",
        "Lack of intelligent automation — existing tools rely on static rules without contextual understanding",
        "Heavy dependence on skilled analysts for threat triage, creating bottlenecks and single points of failure",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.0),
                    problems, font_size=17, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 5: Existing System
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Existing Security Operations", 5)

    existing = [
        "Traditional SIEM tools (Splunk, QRadar, ArcSight) for log aggregation and correlation",
        "Manual triaging by SOC analysts who review alerts one-by-one in queue-based workflows",
        "Static, signature-based and rule-based detection that only catches known attack patterns",
        "Limited contextual intelligence — alerts lack root-cause analysis or actionable explanations",
        "Heavy reliance on human expertise for threat classification and incident prioritization",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.0),
                    existing, font_size=17, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 6: Limitations
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Limitations of Existing Systems", 6)

    limitations = [
        "High false positive rate — up to 90% of SIEM alerts can be benign, wasting analyst time",
        "Time-consuming log investigation — manual correlation across multiple data sources",
        "No AI-based contextual explanation — alerts lack why and how information",
        "Limited automation — response actions still require manual execution",
        "Scalability issues — static rules don't adapt to evolving threat landscapes",
        "Expensive licensing — enterprise SIEM solutions have significant cost barriers",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.0),
                    limitations, font_size=17, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 7: Proposed System
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Proposed System — SOC Copilot", 7)

    proposed = [
        "AI-powered log analysis using hybrid ML (Isolation Forest + Random Forest) for multi-layered detection",
        "Intelligent anomaly detection that identifies unknown/zero-day threats beyond static signatures",
        "Automated threat classification into severity levels: Critical, High, Medium, Low",
        "Natural language explanation of alerts — provides root cause, context, and confidence scores",
        "Automated response suggestions with mitigation steps and recommended actions",
        "Fully offline, desktop-based operation — no cloud dependency, complete data privacy",
        "Governance-first design with kill-switch, audit trail, and human-in-the-loop controls",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5),
                    proposed, font_size=16, spacing=Pt(10))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 8: System Requirements
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "System Requirements", 8)

    # Software card
    add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(4.0), Inches(5.2), CARD_BG)
    add_text_box(slide, Inches(1.0), Inches(1.6), Inches(3.6), Inches(0.4),
                 "⚙️  SOFTWARE", font_size=16, color=ACCENT_CYAN, bold=True)
    sw_items = [
        "Python 3.10+",
        "PyQt6 (Desktop UI Framework)",
        "Scikit-learn (ML Models)",
        "Pandas / NumPy (Data Processing)",
        "SQLite (Feedback & Audit Store)",
        "YAML (Configuration)",
        "PyInstaller (Packaging)",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(2.1), Inches(3.6), Inches(4.2),
                    sw_items, font_size=15, spacing=Pt(8))

    # Hardware card
    add_accent_bar(slide, Inches(5.2), Inches(1.5), Inches(4.0), Inches(5.2), CARD_BG)
    add_text_box(slide, Inches(5.4), Inches(1.6), Inches(3.6), Inches(0.4),
                 "🖥️  HARDWARE", font_size=16, color=ACCENT_GREEN, bold=True)
    hw_items = [
        "Minimum 8 GB RAM",
        "Intel Core i5 or equivalent",
        "2 GB free disk space",
        "No internet required (offline)",
        "Windows / macOS / Linux",
    ]
    add_bullet_list(slide, Inches(5.4), Inches(2.1), Inches(3.6), Inches(4.2),
                    hw_items, font_size=15, spacing=Pt(8), bullet_color=ACCENT_GREEN)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 9: System Architecture
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "System Architecture", 9)

    # Architecture layers
    layers = [
        ("PRESENTATION LAYER", "PyQt6 Desktop UI  •  Dashboard Zones A–F  •  Real-time Alerts Table", ACCENT_CYAN, 1.5),
        ("APPLICATION LAYER", "App Controller  •  Pipeline Orchestrator  •  Governance Engine", ACCENT_BLUE, 2.7),
        ("ML / AI ENGINE", "Isolation Forest (Anomaly)  •  Random Forest (Classification)  •  Ensemble Controller", ACCENT_GREEN, 3.9),
        ("DATA LAYER", "Log Ingestion  •  Preprocessing  •  Feature Engineering  •  SQLite Store", ACCENT_ORANGE, 5.1),
    ]

    for name, desc, color, y in layers:
        # Layer card
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(8.4), Inches(1.0), CARD_BG)
        # Left color indicator
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(0.08), Inches(1.0), color)
        # Layer name
        add_text_box(slide, Inches(1.1), Inches(y + 0.05), Inches(3.0), Inches(0.4),
                     name, font_size=14, color=color, bold=True)
        # Description
        add_text_box(slide, Inches(1.1), Inches(y + 0.45), Inches(7.8), Inches(0.4),
                     desc, font_size=13, color=LIGHT_GRAY)

    # Arrows between layers
    for y in [2.5, 3.7, 4.9]:
        add_text_box(slide, Inches(4.5), Inches(y), Inches(1), Inches(0.3),
                     "⬇", font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 10: Data Flow Diagram
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Data Flow Diagram", 10)

    # Flow steps
    flow_steps = [
        ("01", "LOG INPUT", "User uploads log files or system logs are ingested automatically", ACCENT_CYAN),
        ("02", "PARSING", "Log parser extracts structured fields — timestamp, source IP, event type, payload", ACCENT_BLUE),
        ("03", "PREPROCESSING", "Feature engineering: statistical features, encoding, normalization", ACCENT_GREEN),
        ("04", "ML ANALYSIS", "Isolation Forest detects anomalies → Random Forest classifies attack type", ACCENT_ORANGE),
        ("05", "ALERT ENGINE", "Threats classified by severity (Critical/High/Medium/Low) with confidence scores", ACCENT_RED),
        ("06", "DASHBOARD", "Real-time visualization: threat banner, metric cards, alerts timeline, details", ACCENT_CYAN),
    ]

    for i, (num, title, desc, color) in enumerate(flow_steps):
        y = 1.5 + i * 0.9
        # Number badge
        add_text_box(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.4),
                     num, font_size=20, color=color, bold=True)
        # Step title
        add_text_box(slide, Inches(1.5), Inches(y - 0.02), Inches(2.0), Inches(0.35),
                     title, font_size=14, color=color, bold=True)
        # Description
        add_text_box(slide, Inches(3.5), Inches(y), Inches(5.8), Inches(0.35),
                     desc, font_size=13, color=LIGHT_GRAY)
        # Connector line
        if i < len(flow_steps) - 1:
            add_accent_bar(slide, Inches(1.05), Inches(y + 0.42), Inches(0.02), Inches(0.45), MID_GRAY)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 11: UML Diagrams
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "UML Diagrams", 11)

    # Use Case
    add_accent_bar(slide, Inches(0.5), Inches(1.5), Inches(2.8), Inches(5.2), CARD_BG)
    add_text_box(slide, Inches(0.7), Inches(1.6), Inches(2.4), Inches(0.4),
                 "Use Case Diagram", font_size=14, color=ACCENT_CYAN, bold=True)
    uc_items = [
        "SOC Analyst ↔ Upload Logs",
        "SOC Analyst ↔ View Dashboard",
        "SOC Analyst ↔ Review Alerts",
        "SOC Analyst ↔ Analyze Threats",
        "System ↔ Auto-Ingest Logs",
        "System ↔ Generate Alerts",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.1), Inches(2.4), Inches(4.2),
                    uc_items, font_size=12, spacing=Pt(6))

    # Class Diagram
    add_accent_bar(slide, Inches(3.6), Inches(1.5), Inches(2.8), Inches(5.2), CARD_BG)
    add_text_box(slide, Inches(3.8), Inches(1.6), Inches(2.4), Inches(0.4),
                 "Class Diagram", font_size=14, color=ACCENT_GREEN, bold=True)
    cd_items = [
        "LogParser",
        "FeatureEngineer",
        "IsolationForestModel",
        "RandomForestModel",
        "EnsembleController",
        "AlertEngine",
        "AppController",
        "DashboardUI",
    ]
    add_bullet_list(slide, Inches(3.8), Inches(2.1), Inches(2.4), Inches(4.2),
                    cd_items, font_size=12, spacing=Pt(6), bullet_color=ACCENT_GREEN)

    # Sequence Diagram
    add_accent_bar(slide, Inches(6.7), Inches(1.5), Inches(2.8), Inches(5.2), CARD_BG)
    add_text_box(slide, Inches(6.9), Inches(1.6), Inches(2.4), Inches(0.4),
                 "Sequence Diagram", font_size=14, color=ACCENT_ORANGE, bold=True)
    seq_items = [
        "User → UI: Upload File",
        "UI → Controller: Process",
        "Controller → Parser: Parse",
        "Parser → Features: Extract",
        "Features → ML: Predict",
        "ML → Alerts: Classify",
        "Alerts → UI: Display",
    ]
    add_bullet_list(slide, Inches(6.9), Inches(2.1), Inches(2.4), Inches(4.2),
                    seq_items, font_size=12, spacing=Pt(6), bullet_color=ACCENT_ORANGE)

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 12: Core Modules
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Core Modules", 12)

    modules_core = [
        ("Log Ingestion Module",
         "Supports multiple log formats (CSV, JSON, JSONL, syslog). Implements micro-batch processing for real-time ingestion with configurable batch sizes and intervals."),
        ("Log Parsing & Preprocessing",
         "Extracts structured fields from raw logs. Performs data cleaning, normalization, and statistical feature engineering — including frequency analysis, entropy calculation, and temporal patterns."),
        ("AI Threat Analysis Engine",
         "Hybrid ML approach: Isolation Forest for unsupervised anomaly detection (catches unknown threats) + Random Forest for supervised classification (identifies specific attack types like DDoS, brute-force, infiltration)."),
    ]

    y = 1.5
    for title, desc in modules_core:
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(8.4), Inches(1.6), CARD_BG)
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(0.06), Inches(1.6), ACCENT_BLUE)
        add_text_box(slide, Inches(1.1), Inches(y + 0.1), Inches(8.0), Inches(0.35),
                     title, font_size=18, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, Inches(1.1), Inches(y + 0.5), Inches(8.0), Inches(1.0),
                     desc, font_size=14, color=LIGHT_GRAY)
        y += 1.85

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 13: Support Modules
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "System Support Modules", 13)

    modules_support = [
        ("Alert Classification Module",
         "Classifies threats into Critical, High, Medium, Low severity levels. Assigns confidence scores and generates structured alert objects with full context."),
        ("Recommendation Engine",
         "Provides automated mitigation suggestions based on attack classification. Includes response recommendations, investigation steps, and remediation guidance."),
        ("Dashboard & Visualization",
         "PyQt6 desktop application with zone-based architecture (A–F). Real-time threat banner, system status strip, metric cards, quick actions, and alerts timeline."),
        ("Governance & Audit Module",
         "Kill-switch control, human-in-the-loop enforcement, and audit trail logging. Ensures safety constraints are always respected."),
    ]

    y = 1.5
    for title, desc in modules_support:
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(8.4), Inches(1.2), CARD_BG)
        add_accent_bar(slide, Inches(0.8), Inches(y), Inches(0.06), Inches(1.2), ACCENT_GREEN)
        add_text_box(slide, Inches(1.1), Inches(y + 0.05), Inches(8.0), Inches(0.3),
                     title, font_size=16, color=ACCENT_GREEN, bold=True)
        add_text_box(slide, Inches(1.1), Inches(y + 0.38), Inches(8.0), Inches(0.75),
                     desc, font_size=13, color=LIGHT_GRAY)
        y += 1.4

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 14: Output – Dashboard
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Dashboard Interface", 14, "Output Screens — Overview")

    dashboard_features = [
        ("Zone A — Threat Level Banner",
         "Primary visual indicator showing current threat level (CLEAR / ELEVATED / CRITICAL) with dynamic color coding"),
        ("Zone B — System Status Strip",
         "Consolidated status for Pipeline, Ingestion, and Governance modules with color-coded indicators"),
        ("Zone C — Metric Cards",
         "Clickable cards displaying Total Alerts, Critical, High, Medium, and Low counts with real-time updates"),
        ("Zone D — Quick Actions Bar",
         "Upload files, refresh data, start/stop monitoring — one-click operations for analyst efficiency"),
        ("Zone E — Recent Alerts Timeline",
         "Live feed of latest 10 alerts with priority coloring, timestamps, and clickable details"),
        ("Zone F — Alert Detail Panel",
         "Full alert analysis with severity, classification, log context, and recommended actions"),
    ]

    y = 1.5
    for title, desc in dashboard_features:
        add_text_box(slide, Inches(0.8), Inches(y), Inches(3.2), Inches(0.3),
                     title, font_size=13, color=ACCENT_CYAN, bold=True)
        add_text_box(slide, Inches(4.0), Inches(y), Inches(5.5), Inches(0.35),
                     desc, font_size=12, color=LIGHT_GRAY)
        y += 0.85

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 15: Output – AI Explanation
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "AI-Generated Insights", 15, "Output Screens — Threat Analysis")

    ai_features = [
        "Natural Language Explanation — AI provides human-readable description of detected threats, making complex technical findings accessible to all analysts",
        "Root Cause Analysis — Identifies contributing factors: unusual traffic patterns, anomalous login behavior, suspicious process execution",
        "Recommended Mitigation Steps — Actionable response guidance: block IP, reset credentials, isolate host, escalate to Tier 2",
        "Confidence Score — ML model confidence percentage indicating detection reliability (Isolation Forest anomaly score + Random Forest probability)",
        "Attack Classification — Specific attack type identification: DDoS, Brute Force, Port Scan, Infiltration, Botnet, with supporting evidence",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5),
                    ai_features, font_size=15, spacing=Pt(14))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 16: Conclusion
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Conclusion", 16)

    conclusions = [
        "SOC Copilot significantly reduces SOC analyst workload by automating repetitive log analysis and alert triage tasks",
        "Enables faster incident response through intelligent threat detection and automated classification — reducing MTTD and MTTR",
        "Hybrid ML approach (Isolation Forest + Random Forest) provides comprehensive threat coverage for both known and unknown attacks",
        "Governance-first design with kill-switch and audit trails ensures safe, accountable AI-driven security operations",
        "Fully offline, scalable desktop application that respects data privacy and works without cloud dependencies",
        "Represents a practical, deployable AI assistant that augments — not replaces — human security analysts",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5),
                    conclusions, font_size=16, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 17: Future Scope
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)
    add_header_with_bar(slide, "Future Enhancements", 17)

    future = [
        "Real-time SIEM Integration — Connect with Splunk, Elastic SIEM, and QRadar for live log streaming and bi-directional alert sync",
        "Advanced Deep Learning — Autoencoder and Transformer-based models for more sophisticated anomaly detection and sequential attack pattern recognition",
        "Cloud Deployment — Containerized deployment with Docker/Kubernetes for enterprise-scale SOC environments",
        "Automated Remediation Scripts — Execute pre-approved response actions (firewall rules, account lockouts) with governance controls",
        "Multi-Tenant Enterprise Version — Role-based access control, team dashboards, and centralized management for large SOC teams",
        "Threat Intelligence Feeds — Integration with MITRE ATT&CK, VirusTotal, and other threat intelligence platforms for enriched context",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5),
                    future, font_size=15, spacing=Pt(12))

    # ═══════════════════════════════════════════════════════════════
    # SLIDE 18: Thank You
    # ═══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    # Decorative elements
    add_accent_bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.08), ACCENT_BLUE)
    add_accent_bar(slide, Inches(0), Inches(7.42), Inches(10), Inches(0.08), ACCENT_CYAN)

    # Thank you text
    add_text_box(slide, Inches(0), Inches(2.0), Inches(10), Inches(1.0),
                 "Thank You", font_size=52, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(3.0), Inches(10), Inches(0.5),
                 "SOC Copilot — AI-Powered Security Operations Assistant",
                 font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

    # Divider
    add_accent_bar(slide, Inches(3.5), Inches(3.7), Inches(3), Inches(0.03), MID_GRAY)

    # Team info
    add_text_box(slide, Inches(0), Inches(4.0), Inches(10), Inches(0.4),
                 "Team Members", font_size=16, color=LIGHT_GRAY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0), Inches(4.4), Inches(10), Inches(0.4),
                 "[Member 1]  •  [Member 2]  •  [Member 3]  •  [Member 4]",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(5.0), Inches(10), Inches(0.4),
                 "Contact: [email@institution.edu]",
                 font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Q&A
    add_text_box(slide, Inches(0), Inches(5.7), Inches(10), Inches(0.5),
                 "Questions & Answers",
                 font_size=20, color=ACCENT_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 18)

    # ── Save ──────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__), "SOC_Copilot_Presentation.pptx")
    prs.save(output_path)
    print(f"[OK] Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
