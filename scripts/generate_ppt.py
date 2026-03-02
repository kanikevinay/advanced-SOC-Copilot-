"""
SOC Copilot Review-1 PowerPoint Generator
Generates a professional PPTX presentation with diagrams and visuals.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Color scheme - Professional Cybersecurity theme
COLORS = {
    'primary': RGBColor(0, 112, 192),      # Blue
    'secondary': RGBColor(0, 176, 80),     # Green
    'accent': RGBColor(255, 192, 0),       # Orange/Gold
    'dark': RGBColor(31, 73, 125),         # Dark Blue
    'critical': RGBColor(192, 0, 0),       # Red
    'text': RGBColor(51, 51, 51),          # Dark Gray
    'white': RGBColor(255, 255, 255),
    'light_bg': RGBColor(242, 242, 242),
}


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['dark']
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['dark']
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(12)
    
    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['dark']
    title_shape.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Left column header
    left_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.9), Inches(0.6)
    )
    left_header.fill.solid()
    left_header.fill.fore_color.rgb = COLORS['primary']
    left_header.line.fill.background()
    left_header.text_frame.paragraphs[0].text = left_title
    left_header.text_frame.paragraphs[0].font.size = Pt(18)
    left_header.text_frame.paragraphs[0].font.bold = True
    left_header.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    left_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5.9), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(8)
    
    # Right column header
    right_header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(5.9), Inches(0.6)
    )
    right_header.fill.solid()
    right_header.fill.fore_color.rgb = COLORS['secondary']
    right_header.line.fill.background()
    right_header.text_frame.paragraphs[0].text = right_title
    right_header.text_frame.paragraphs[0].font.size = Pt(18)
    right_header.text_frame.paragraphs[0].font.bold = True
    right_header.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    right_header.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(2.2), Inches(5.9), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(8)
    
    return slide


def add_architecture_slide(prs):
    """Add the system architecture diagram slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['dark']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SYSTEM ARCHITECTURE"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Main container
    main_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.0)
    )
    main_box.fill.solid()
    main_box.fill.fore_color.rgb = COLORS['light_bg']
    main_box.line.color.rgb = COLORS['dark']
    
    # Phase 1 - Detection Engine
    phase1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.0)
    )
    phase1.fill.solid()
    phase1.fill.fore_color.rgb = COLORS['primary']
    phase1.text_frame.paragraphs[0].text = "PHASE-1: DETECTION ENGINE"
    phase1.text_frame.paragraphs[0].font.size = Pt(14)
    phase1.text_frame.paragraphs[0].font.bold = True
    phase1.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    
    # Pipeline boxes inside Phase 1
    boxes = ["Log\nIngestion", "Preprocessing", "Feature\nEngineering", "Isolation\nForest", "Random\nForest", "Ensemble\nCoordinator"]
    x_positions = [0.7, 2.5, 4.3, 6.1, 8.1, 10.2]
    
    for i, (box_text, x) in enumerate(zip(boxes, x_positions)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(1.6), Inches(1.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = COLORS['primary']
        tf = box.text_frame
        tf.paragraphs[0].text = box_text
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
        
        # Add arrows
        if i < len(boxes) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.65), Inches(2.4), Inches(0.3), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['accent']
            arrow.line.fill.background()
    
    # Phase 2 - Trust & Intelligence
    phase2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.6), Inches(12.3), Inches(1.4)
    )
    phase2.fill.solid()
    phase2.fill.fore_color.rgb = COLORS['secondary']
    phase2.text_frame.paragraphs[0].text = "PHASE-2: TRUST & INTELLIGENCE"
    phase2.text_frame.paragraphs[0].font.size = Pt(14)
    phase2.text_frame.paragraphs[0].font.bold = True
    phase2.text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    
    # Phase 2 components
    phase2_boxes = ["Feedback Store", "Drift Monitoring", "Threshold Calibration", "Explainability"]
    phase2_x = [1.0, 4.0, 7.0, 10.0]
    
    for box_text, x in zip(phase2_boxes, phase2_x):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.1), Inches(2.5), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = COLORS['secondary']
        tf = box.text_frame
        tf.paragraphs[0].text = box_text
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
    
    # Phase 3 - Governance
    phase3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.4)
    )
    phase3.fill.solid()
    phase3.fill.fore_color.rgb = COLORS['accent']
    phase3.text_frame.paragraphs[0].text = "PHASE-3: GOVERNANCE INFRASTRUCTURE (Disabled by Default)"
    phase3.text_frame.paragraphs[0].font.size = Pt(14)
    phase3.text_frame.paragraphs[0].font.bold = True
    phase3.text_frame.paragraphs[0].font.color.rgb = COLORS['dark']
    
    # Phase 3 components
    phase3_boxes = ["Governance Policy", "Approval Workflow", "Kill Switch", "Audit Logger"]
    phase3_x = [1.0, 4.0, 7.0, 10.0]
    
    for box_text, x in zip(phase3_boxes, phase3_x):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.7), Inches(2.5), Inches(0.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['white']
        box.line.color.rgb = COLORS['accent']
        tf = box.text_frame
        tf.paragraphs[0].text = box_text
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = COLORS['dark']
    
    return slide


def add_algorithm_slide(prs):
    """Add algorithms used slide with visual boxes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['dark']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ALGORITHMS USED - HYBRID ML APPROACH"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Isolation Forest box
    if_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6.1), Inches(2.8)
    )
    if_box.fill.solid()
    if_box.fill.fore_color.rgb = COLORS['primary']
    if_box.line.fill.background()
    
    # IF header
    if_header = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.7), Inches(0.5))
    tf = if_header.text_frame
    p = tf.paragraphs[0]
    p.text = "🌲 ISOLATION FOREST"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # IF content
    if_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.7), Inches(2.0))
    tf = if_content.text_frame
    tf.word_wrap = True
    items = [
        "Type: Unsupervised Anomaly Detection",
        "Training: Benign records only",
        "Output: Anomaly score [0, 1]",
        "Purpose: Detect novel, unknown threats"
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(6)
    
    # Random Forest box
    rf_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6.1), Inches(2.8)
    )
    rf_box.fill.solid()
    rf_box.fill.fore_color.rgb = COLORS['secondary']
    rf_box.line.fill.background()
    
    # RF header
    rf_header = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.5))
    tf = rf_header.text_frame
    p = tf.paragraphs[0]
    p.text = "🌳 RANDOM FOREST"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # RF content
    rf_content = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.7), Inches(2.0))
    tf = rf_content.text_frame
    tf.word_wrap = True
    items = [
        "Type: Supervised Classification",
        "Training: Full labeled dataset",
        "Output: Class + confidence",
        "Classes: DDoS, Malware, BruteForce..."
    ]
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['white']
        p.space_after = Pt(6)
    
    # Ensemble box at bottom
    ensemble_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.4), Inches(8.3), Inches(1.0)
    )
    ensemble_box.fill.solid()
    ensemble_box.fill.fore_color.rgb = COLORS['accent']
    ensemble_box.line.fill.background()
    
    ensemble_text = slide.shapes.add_textbox(Inches(2.7), Inches(4.6), Inches(7.9), Inches(0.6))
    tf = ensemble_text.text_frame
    p = tf.paragraphs[0]
    p.text = "⚡ ENSEMBLE COORDINATOR → Risk Score + Priority (P0-P4)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER
    
    # Arrows pointing to ensemble
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(3.5), Inches(4.15), Inches(0.3), Inches(0.25)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLORS['primary']
    arrow1.line.fill.background()
    
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(9.5), Inches(4.15), Inches(0.3), Inches(0.25)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLORS['secondary']
    arrow2.line.fill.background()
    
    # Why Hybrid section
    why_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.6), Inches(12.4), Inches(1.5)
    )
    why_box.fill.solid()
    why_box.fill.fore_color.rgb = COLORS['light_bg']
    why_box.line.color.rgb = COLORS['dark']
    
    why_header = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.4))
    tf = why_header.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 Why Hybrid Approach?"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    why_content = slide.shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.9))
    tf = why_content.text_frame
    items = ["Unknown threats detected by Isolation Forest", "Known threats classified by Random Forest", 
             "Two signals provide higher confidence", "Robustness: backup if one model fails"]
    p = tf.paragraphs[0]
    p.text = "  ✓ " + "    ✓ ".join(items)
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text']
    
    return slide


def add_tools_slide(prs):
    """Add tools and technologies slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.0)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS['dark']
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TOOLS & TECHNOLOGIES"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # Four quadrant layout
    categories = [
        ("Frontend", COLORS['primary'], ["PyQt6 - Desktop GUI", "CLI Interface", "Real-time Dashboard"]),
        ("Backend", COLORS['secondary'], ["Python 3.10+", "Scikit-learn ML", "Pandas/NumPy"]),
        ("Database", COLORS['accent'], ["SQLite - Feedback", "YAML - Config", "Joblib - Models"]),
        ("Dev Tools", COLORS['critical'], ["pytest - Testing", "Black/Ruff - Linting", "Git - Version Control"])
    ]
    
    positions = [(0.5, 1.3), (6.8, 1.3), (0.5, 4.3), (6.8, 4.3)]
    
    for (cat_name, color, items), (x, y) in zip(categories, positions):
        # Category box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.0), Inches(2.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Header
        header = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(5.6), Inches(0.5))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = cat_name
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        
        # Items
        content = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(5.6), Inches(1.8))
        tf = content.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['white']
            p.space_after = Pt(8)
    
    return slide


def add_thank_you_slide(prs, team_members):
    """Add thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['dark']
    shape.line.fill.background()
    
    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.333), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SOC Copilot - Making Security Operations Smarter, Safer, and More Transparent"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Team members box
    team_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.0), Inches(6.333), Inches(2.5)
    )
    team_box.fill.solid()
    team_box.fill.fore_color.rgb = RGBColor(50, 90, 140)
    team_box.line.fill.background()
    
    # Team header
    team_header = slide.shapes.add_textbox(Inches(3.7), Inches(4.1), Inches(5.9), Inches(0.5))
    tf = team_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Team"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER
    
    # Team content
    team_content = slide.shapes.add_textbox(Inches(3.7), Inches(4.6), Inches(5.9), Inches(1.8))
    tf = team_content.text_frame
    for i, (name, rno) in enumerate(team_members):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {name} - {rno}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
    
    # Questions
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def generate_presentation():
    """Generate the complete SOC Copilot presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "SOC COPILOT", "Hybrid ML-Based Security Operations Center Assistant\nReview-1 Presentation")
    
    # Slide 2: Overview
    add_content_slide(prs, "OVERVIEW", [
        "Problem Statement: SOC analysts face alert fatigue and delayed threat response",
        "Objective: Build fully offline, hybrid ML-based security assistant",
        "Review-1 Focus: Requirement analysis, system design, module identification",
        "Modules: Log Ingestion, ML Detection, Alert Generation, Governance",
        "Algorithms: Isolation Forest + Random Forest (Hybrid Ensemble)",
        "Technologies: Python, Scikit-learn, PyQt6, SQLite"
    ])
    
    # Slide 3: Problem Statement
    add_content_slide(prs, "PROBLEM STATEMENT", [
        "Alert Fatigue: Analysts review thousands of alerts daily, missing genuine threats",
        "Delayed Response: Manual triage increases Mean Time to Detect (MTTD)",
        "Inconsistent Decisions: Human judgment varies across analysts and shifts",
        "Data Overload: Unable to process heterogeneous log formats efficiently",
        "Lack of Explainability: Existing tools provide opaque, non-actionable alerts",
        "Cloud Dependency: Most SIEM solutions unsuitable for air-gapped environments"
    ])
    
    # Slide 4: Objectives
    add_content_slide(prs, "OBJECTIVES", [
        "Threat Detection: Identify malicious activity using hybrid ML (unsupervised + supervised)",
        "Alert Prioritization: Generate P0-P4 priority alerts to focus analyst attention",
        "Explainability: Provide human-readable reasoning for every decision",
        "Offline Operation: Operate without network connectivity or cloud dependencies",
        "Governance-First: Ensure all automation is disabled by default with manual controls",
        "Performance: Achieve >80% classification accuracy with minimal false positives"
    ])
    
    # Slide 5: Architecture
    add_architecture_slide(prs)
    
    # Slide 6: Module Description Part 1
    add_two_column_slide(prs, "MODULE DESCRIPTION (Part 1)",
        "Module 1: Log Ingestion",
        ["Multi-format: JSON, CSV, Syslog, EVTX", "Format detection & field extraction",
         "Timestamp normalization to UTC", "Schema & field validation"],
        "Module 2: Data Preprocessing",
        ["Field standardization (src_ip, dst_ip, port)",
         "Categorical to integer encoding", "Missing value handling",
         "Normalized DataFrame output"])
    
    # Slide 7: Module Description Part 2
    add_two_column_slide(prs, "MODULE DESCRIPTION (Part 2)",
        "Module 3: ML Detection",
        ["78 numeric feature extraction", "Isolation Forest anomaly scores",
         "Random Forest classification", "Ensemble risk scoring"],
        "Module 4: Alert & Reporting",
        ["Priority assignment (P0-P4)", "MITRE ATT&CK mapping",
         "Feature importance explanation", "Real-time dashboard"])
    
    # Slide 8: Algorithms
    add_algorithm_slide(prs)
    
    # Slide 9: Techniques
    add_content_slide(prs, "TECHNIQUES USED", [
        "Data Preprocessing: Timestamp normalization, categorical encoding, field standardization",
        "Feature Engineering: 78 features (statistical, temporal, behavioral, network)",
        "ML Optimization: Class imbalance handling, weighted ensemble, decision matrix",
        "Security: Offline-first, read-only models, deterministic scoring",
        "Explainability: Wrapper-based feature importance, post-hoc reasoning",
        "Governance: Append-only audit logging, kill switch, approval workflows"
    ])
    
    # Slide 10: Tools
    add_tools_slide(prs)
    
    # Slide 11: Review Outcome
    add_content_slide(prs, "REVIEW-1 OUTCOME", [
        "✅ Requirements Finalized: Offline, hybrid ML, multi-format ingestion",
        "✅ Modules Identified: Log Ingestion, Preprocessing, ML Detection, Alerts, Governance",
        "✅ Architecture Designed: Three-phase modular system",
        "✅ Algorithms Selected: Isolation Forest + Random Forest ensemble",
        "✅ Tech Stack Finalized: Python, Scikit-learn, PyQt6, SQLite",
        "✅ Ready for Implementation Phase with 208+ test cases planned"
    ])
    
    # Slide 12: Thank You
    team_members = [
        ("Member 1", "RNO"),
        ("Member 2", "RNO"),
        ("Member 3", "RNO"),
        ("Member 4", "RNO"),
    ]
    add_thank_you_slide(prs, team_members)
    
    # Save
    output_path = Path(__file__).parent.parent / "SOC_Copilot_Review1_Presentation.pptx"
    prs.save(str(output_path))
    print(f"[SUCCESS] Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_presentation()
